from langchain_community.document_loaders import PyMuPDFLoader, UnstructuredPowerPointLoader
from langchain_core.document_loaders import BaseLoader
from langchain_core.documents import Document
import os 
import uuid
import requests
import mimetypes
from typing import List
from urllib.parse import urljoin
import os
import re
import fitz
from pptx import Presentation
from dotenv import load_dotenv
import time
from pathlib import Path
load_dotenv()


"""Custom URL loader usig jina reader ai"""
class UrlLoader(BaseLoader):
    def __init__(self,url):
        self.api_key = os.getenv("JINA_API_KEY")
        self.url = url
        self.jina_url = f"https://r.jina.ai/{url}"
    
    def lazy_load(self):
        headers = {
            "Authorization": f"Bearer {self.api_key}",
            "X-Retain-Images": "all",
            "X-With-Shadow-Dom": "true",
            "X-With-Iframe": "true",
            "X-Cache-Tolerance": "3600"
        }
        try:
            response = requests.get(url=self.jina_url,headers=headers,timeout=15)
            response.raise_for_status()

            markdown_content = response.text

            yield Document(
                page_content=markdown_content,
                metadata = {
                    "source" : self.url
                }
            )
        except requests.exceptions.RequestException as e:
            print(f"Error in fetching url: {self.url}")

    def get_images(self,markdown_content,save_dir = r".\temp-images"):
        os.makedirs(save_dir,exist_ok=True)
        images_urls = re.findall(r'!\[.*?\]\((.*?)\)',markdown_content)

        local_img_paths = []

        for img_url in images_urls:
            abs_url = urljoin(self.url,img_url)
            
            try:
                img_reponse = requests.get(url=abs_url,stream=True,timeout=15)
                img_reponse.raise_for_status()

                content_type = img_reponse.headers.get('content-type')
                ext = mimetypes.guess_extension(content_type) or '.jpg'

                filename = f"{uuid.uuid4()}{ext}"
                local_path = os.path.join(save_dir,filename)

                with open(local_path,"wb") as img:
                    for chunk in img_reponse.iter_content(8192):
                        img.write(chunk)

                local_img_paths.append({
                    "image_path": local_path,
                    "metadata": {
                        "source": abs_url,
                        "type": "image"
                    }
                })
            except requests.exceptions.RequestException :
                print(f"File download unsuccessfull. {abs_url}")
        return local_img_paths

class Loaders():
    def __init__(self,sources:List[str],save_dir = r".\temp-images"):
        self.sources = sources
        self.save_dir = save_dir
        """Dividing data depending on the file extension"""
        print(f"Reading sources and separating them....")
        time.sleep(1)
        self.pdfs = []
        self.ppts = []
        self.urls = []
        for src in sources:
            src = src.lower()
            if(src.startswith(("https://","http://"))):
                self.urls.append(src)
            elif(src.endswith(".pdf")):
                self.pdfs.append(src)
            elif(src.endswith(".pptx")):
                self.ppts.append(src)

        self.all_doc = []
        self.imgs = []

        if self.pdfs:
            print(f"Gathering PDFs..")
            time.sleep(1)
            self.pdf_loader()
        if self.ppts:
            print(f"Gathering PPTs..")
            time.sleep(1)
            self.ppt_loader()
        if self.urls:
            print(f"Gathering URLs..")
            time.sleep(1)
            self.url_loader()

        self.convert_all_svgs_to_png(self.save_dir)
    
    def convert_all_svgs_to_png(self,image_dir=r".\temp-images"):
        """Converting SVGs to PNGs"""
        image_dir = Path(image_dir)
        svg_files = list(image_dir.glob("**/*.svg"))
        
        print(f"Found {len(svg_files)} SVG files to convert...")
        
        for svg_path in svg_files:
            png_path = svg_path.with_suffix(".png")
            doc = fitz.open(svg_path)
            page = doc[0]
            pix = page.get_pixmap()
            pix.save(str(png_path))
            doc.close()
            print(f"Converted: {svg_path.name} -> {png_path.name}")
            svg_path.unlink()
            
    
    def get_pdf_sync_images(self,file_path,):
        print(f"Gathering images in PDF {file_path}")
        os.makedirs(self.save_dir,exist_ok=True)
        source_document = fitz.open(file_path)
        for page_num in range(len(source_document)):
            page = source_document[page_num]
            image_list = page.get_images(full=True)
            for img in image_list:
                xref = img[0]
                base_image = source_document.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                filename = f"pdf_{uuid.uuid4()}.{image_ext}"
                filepath = os.path.join(self.save_dir, filename)
                with open(filepath, "wb") as f:
                    f.write(image_bytes)
                self.imgs.append({
                    "image_path": filepath,
                    "metadata": {
                        "source": file_path,
                        "page": page_num,
                        "type": "image"
                    }
                })

    def get_ppt_sync_images(self,file_path):
        print(f"Gathering images in PPT {file_path}")
        os.makedirs(self.save_dir,exist_ok=True)
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "image"): 
                    image = shape.image
                    image_bytes = image.blob
                    image_ext = image.ext
                    filename = f"ppt_{uuid.uuid4()}.{image_ext}"
                    filepath = os.path.join(self.save_dir, filename)
                    with open(filepath, "wb") as f:
                        f.write(image_bytes)
                    self.imgs.append({
                        "image_path": filepath,
                        "metadata": {
                            "source": file_path,
                            "page": slide_num,
                            "type": "image"
                        }
                    })

    def pdf_loader(self):
        for _ in self.pdfs:
            loader = PyMuPDFLoader(_)
            doc = loader.load()
            self.get_pdf_sync_images(_)
            self.all_doc.extend(doc)

    def ppt_loader(self):
        for _ in self.ppts:
            loader = UnstructuredPowerPointLoader(_)
            doc = loader.load()
            self.get_ppt_sync_images(_)
            self.all_doc.extend(doc)

    def url_loader(self):
        for _ in self.urls:
            loader = UrlLoader(_)
            doc = list(loader.lazy_load())
            for i in range(len(doc)):
                markdown_content = doc[i].page_content
                print(f"Gathering images in URLs..")
                self.imgs.extend(loader.get_images(markdown_content=markdown_content))
            self.all_doc.extend(doc)


# loader = Loaders(sources)
# print(loader.imgs)
# for i in loader.imgs:
#     print(i["metadata"])
        
# loader = UrlLoader("https://lilianweng.github.io/posts/2018-06-24-attention/")
# docs = list(loader.lazy_load())
# if docs:
#     print(docs)
#     markdown_content = docs[0].page_content
#     imgs = loader.get_images(markdown_content=markdown_content)
#     print(imgs)
#     print(docs)
#     print(docs[0].page_content[:500])    
